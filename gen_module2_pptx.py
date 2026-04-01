#!/usr/bin/env python3
"""
NLB Prishtina — Module 2: Product Mix, Funding & POS Strategy
4-Slide Executive PPTX | McKinsey-style
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x00,0x2B,0x49)
STEEL = RGBColor(0x4A,0x6F,0xA5)
CHARCOAL = RGBColor(0x1A,0x1A,0x2E)
WHITE = RGBColor(0xFF,0xFF,0xFF)
GREEN = RGBColor(0x05,0x96,0x69)
RED = RGBColor(0xDC,0x26,0x26)
AMBER = RGBColor(0xD9,0x77,0x06)
TEAL = RGBColor(0x0D,0x94,0x88)
GRAY = RGBColor(0x64,0x74,0x8B)
BORDER = RGBColor(0xE2,0xE8,0xF0)
LIGHT = RGBColor(0xF0,0xF4,0xF8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

with open('branch_data_unified.json') as f:
    data = json.load(f)
branches = data['branches']

POS = {'0001':1705,'1006':421,'100104':201,'1005':156,'100106':126,'1007':107,
       '1002':115,'1003':107,'100603':71,'100703':81,'100702':69,'000103':78,
       '100201':50,'100602':62,'100701':47,'1009':45,'1004':60,'100503':56,
       '100501':43,'100604':54,'100302':42,'100502':33,'100301':27,'100202':28,
       '100706':6,'100105':11,'100705':3,'100303':9,'100108':6}

MCLICK = {'100706':75.0,'0001':62.4,'100301':56.2,'100104':56.0,'100604':51.8,
          '100503':50.2,'100302':49.6,'000103':49.5,'1002':48.3,'100603':48.2,
          '100705':48.1,'100602':46.5,'1003':46.3,'100202':45.8,'100201':45.4,
          '100105':45.1,'1009':43.7,'1006':43.6,'100106':41.5,'100502':38.5,
          '100702':37.6,'1004':37.6,'1005':37.0,'1007':35.7,'100701':34.9,
          '100501':33.8,'100703':30.8,'100108':28.8,'100303':27.2}

def s(v,d=0):
    if v is None: return d
    try: return float(v)
    except: return d

# Build branch data
BB = []
for ou, b in branches.items():
    if not b.get('is_benchmarked'): continue
    emp = b['employees'] or 1
    pi_dep = s(b.get('pi_deposits'))
    ca = s(b.get('pi_ca_balance'))
    tda = s(b.get('pi_tda_balance'))
    sa = s(b.get('pi_savings_balance'))
    pi_tot = pi_dep if pi_dep > 0 else 1
    active = s(b.get('active_clients_pi')) or 1
    cl = s(b.get('consumer_loan_no'))
    cl_pen = cl / active * 100
    pos = POS.get(ou, 0)
    biz = s(b.get('businesses')) or 1
    BB.append({
        'ou':ou,'name':b['branch_name'],'city':b.get('city',''),'emp':b['employees'],
        'ca_pct':ca/pi_tot*100,'tda_pct':tda/pi_tot*100,'sa_pct':sa/pi_tot*100,
        'cl':int(cl),'cl_pen':cl_pen,'active':int(active),
        'mort':int(s(b.get('housing_loan_no'))),'cc':int(s(b.get('credit_cards_no'))),
        'ovd':int(s(b.get('overdraft_no'))),
        'pos':pos,'biz':int(biz),'pos_den':pos/biz*100,
        'mclick':MCLICK.get(ou,0),
        'xsell':s(b.get('cross_sell_intensity')),
        'pi_ob':s(b.get('pi_onbalance_segment',b.get('total_pi_onbalance',0))),
        'micro_ob':s(b.get('micro_onbalance_segment',0)),
        'cib_ob':s(b.get('cib_onbalance_segment',0)),
        'total_dep':s(b.get('total_deposits')),
        'onbal':s(b.get('total_onbalance')),
    })

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

def box(sl,l,t,w,h,fill=None,brd=None,bw=Pt(1)):
    shp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if brd: shp.line.color.rgb = brd; shp.line.width = bw
    else: shp.line.fill.background()
    return shp

def txt(sl,l,t,w,h,text,sz=11,bold=False,clr=CHARCOAL,align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = clr; p.alignment = align
    return tf

def exec_bar(sl, y, label, text):
    shp = box(sl, Inches(0.4), y, Inches(12.5), Inches(0.65), fill=LIGHT, brd=BORDER)
    tf = shp.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label + '  '; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = text; r2.font.size = Pt(9); r2.font.color.rgb = CHARCOAL

def header_bar(sl, title, sub):
    box(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.9), fill=NAVY)
    txt(sl, Inches(0.5), Inches(0.12), Inches(10), Inches(0.4), title, sz=20, bold=True, clr=WHITE)
    txt(sl, Inches(0.5), Inches(0.52), Inches(10), Inches(0.3), sub, sz=10, clr=RGBColor(0xA0,0xBC,0xD8))

def kpi_box(sl, x, y, label, value, ctx, accent):
    box(sl, x, y, Inches(3.8), Inches(0.85), fill=WHITE, brd=BORDER)
    box(sl, x, y, Inches(0.06), Inches(0.85), fill=accent)
    txt(sl, x+Inches(0.2), y+Inches(0.08), Inches(3.4), Inches(0.2), label, sz=8, bold=True, clr=GRAY)
    txt(sl, x+Inches(0.2), y+Inches(0.28), Inches(3.4), Inches(0.3), value, sz=22, bold=True, clr=accent)
    txt(sl, x+Inches(0.2), y+Inches(0.62), Inches(3.4), Inches(0.2), ctx, sz=8, clr=GRAY)

# ============================================================
# SLIDE 1 — PRODUCT MIX OVERVIEW
# ============================================================
sl1 = prs.slides.add_slide(blank)
header_bar(sl1, 'MODULE 2 \u2014 PRODUCT MIX & FUNDING STRATEGY', 'NLB Prishtina | Consumer Loans, Deposit Mix, POS, Digital Adoption')

exec_bar(sl1, Inches(1.1), 'EXECUTIVE INSIGHT:',
    'Consumer loan penetration averages 5.7% \u2014 half the network is below. Three branches exceed 50% TDA dependency (Gjakove, Prizren 2, Fushe Kosove). '
    'Five branches have zero POS terminals including 1001 Prishtine serving 21,396 businesses. No pricing/rate data is available in source \u2014 funding cost logic is directional.')

# KPIs
kpi_box(sl1, Inches(0.4), Inches(2.0), 'CONSUMER LOAN PENETRATION', '5.7%', 'Median: consumer loans / active PI clients', NAVY)
kpi_box(sl1, Inches(4.5), Inches(2.0), 'TDA SHARE (PI DEPOSITS)', '36.6%', 'Median. 3 branches above 50%', AMBER)
kpi_box(sl1, Inches(8.6), Inches(2.0), 'POS NETWORK TERMINALS', '3,819', 'Density: 2.6% of ~148K businesses', TEAL)

# Product mix summary table (top issues)
txt(sl1, Inches(0.4), Inches(3.1), Inches(12), Inches(0.3),
    'CRITICAL PRODUCT MIX ISSUES BY BRANCH', sz=11, bold=True, clr=NAVY)

issues = [
    ['Filiali Prizren (1005)','Consumer 3.2% (worst main Filiali)','TDA 35.8% (balanced)','POS 2.2% density','Push consumer hard; POS growth in 7,119-biz market'],
    ['Filiali Besiane (1009)','Consumer 3.6% (bottom 5)','TDA 32.9% (good)','POS 1.9% density','Push consumer + POS; strong CA base (63.6%)'],
    ['Filiali Gjakove (1004)','Consumer 5.3% (below median)','TDA 54.4% (CRITICAL)','POS 1.5% density','Moderate TDA; shift to CA/salary; push POS in 4,005-biz market'],
    ['Eksp. Prizren 2 (100504)','Consumer 6.4% (above median)','TDA 56.8% (CRITICAL)','POS 0.0% (zero)','Deploy POS immediately; moderate TDA growth'],
    ['Eksp. Fushe Kosove (100104)','Consumer 9.8% (strong)','TDA 50.0% (HIGH)','POS 7.4% density','Funding mix correction needed; consumer already strong'],
    ['1001 Prishtine (1001)','Consumer 6.6% (near median)','TDA 41.7% (elevated)','POS 0.0% (zero)','Deploy POS in 21,396-biz market; monitor TDA'],
    ['Filiali Mitrovice (1007)','Consumer 3.9% (bottom 6)','TDA 36.5% (balanced)','POS 4.2% density','Consumer loan campaign priority; product mix too shallow'],
    ['Eksp. Decan (100604)','Consumer 3.5% (bottom 5)','TDA 45.1% (elevated)','POS 6.8% density','Push consumer aggressively; 7,528 active clients under-served'],
]

tbl = sl1.shapes.add_table(len(issues)+1, 5, Inches(0.4), Inches(3.4), Inches(12.5), Inches(3.8))
t = tbl.table
t.columns[0].width = Inches(2.5); t.columns[1].width = Inches(2.5); t.columns[2].width = Inches(2.2)
t.columns[3].width = Inches(2.0); t.columns[4].width = Inches(3.3)

for ci, h in enumerate(['Branch','Consumer Loan Issue','TDA Issue','POS Issue','Management Action']):
    c = t.cell(0,ci); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs: p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = WHITE

for ri, row in enumerate(issues):
    for ci, v in enumerate(row):
        c = t.cell(ri+1,ci); c.text = v
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(7); p.font.color.rgb = CHARCOAL
            if 'CRITICAL' in v: p.font.color.rgb = RED; p.font.bold = True
            elif 'zero' in v.lower(): p.font.color.rgb = RED

notes1 = sl1.notes_slide
notes1.notes_text_frame.text = """SLIDE 1 — PRODUCT MIX OVERVIEW

KEY POINTS:
- Consumer loan penetration median is 5.7%. Prizren (3.2%), Besiane (3.6%), Viti (3.5%), Decan (3.5%) are the worst.
- TDA dependency: Prizren 2 (56.8%), Gjakove (54.4%), Fushe Kosove (50.0%) exceed 50% threshold.
- Five branches have ZERO POS: 1001 Prishtine, Prizren 2, Prishtine-Llapit, Bregu i Diellit, Mitrovice 2.
  1001 Prishtine alone sits in a market with 21,396 businesses and has zero POS terminals — this is the single largest
  merchant-acquiring gap in the network.

PRICING DATA NOTE:
No pricing or rate data exists in any source file. All funding-cost recommendations are based on the structural
hierarchy: TDA is the most expensive PI funding source > Savings > CA. Specific rate spreads cannot be quantified.
If management wants rate-based analysis, pricing data must be provided separately.

IMPORTANT DISTINCTIONS:
- Consumer loan counts are FLOW (2025 new production), not stock
- Consumer loan penetration = flow / active base — measures sales intensity
- TDA/CA/SA percentages are STOCK (current balances), not flows
- Do not confuse the two — one is sales effort, the other is portfolio structure
"""

# ============================================================
# SLIDE 2 — CONSUMER LOAN GROWTH STRATEGY
# ============================================================
sl2 = prs.slides.add_slide(blank)
header_bar(sl2, 'CONSUMER LOAN GROWTH \u2014 BRANCH PRIORITIES', 'Where to push, where to defend, where capacity exists')

exec_bar(sl2, Inches(1.1), 'SO WHAT:',
    'Consumer lending is the highest-margin PI product and the primary commercial lever. 15 branches are below the 5.7% median. '
    'The 5 worst-performing branches (Prizren 3.2%, Decan 3.5%, Viti 3.5%, Besiane 3.6%, Mitrovice 3.9%) together serve 53,000+ active clients \u2014 '
    'a massive addressable pool for consumer loan campaigns.')

# Consumer loan table — all 31 branches sorted by penetration
cl_sorted = sorted(BB, key=lambda x: x['cl_pen'])

# Split: below median vs above
below = [b for b in cl_sorted if b['cl_pen'] < 5.71]
above = [b for b in cl_sorted if b['cl_pen'] >= 5.71]

txt(sl2, Inches(0.4), Inches(1.95), Inches(6), Inches(0.25),
    f'PUSH CONSUMER GROWTH ({len(below)} branches below 5.7% median)', sz=10, bold=True, clr=RED)

# Left table — below median
ltbl = sl2.shapes.add_table(min(len(below)+1,17), 5, Inches(0.4), Inches(2.25), Inches(6.2), Inches(4.5))
lt = ltbl.table
lt.columns[0].width = Inches(2.3); lt.columns[1].width = Inches(1.0); lt.columns[2].width = Inches(0.8)
lt.columns[3].width = Inches(1.0); lt.columns[4].width = Inches(1.1)

for ci,h in enumerate(['Branch','City','CL Pen%','Active PI','CL Count']):
    c = lt.cell(0,ci); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = RED
    for p in c.text_frame.paragraphs: p.font.size = Pt(7); p.font.bold = True; p.font.color.rgb = WHITE

for ri,b in enumerate(below[:16]):
    vals = [b['name'],b['city'],f"{b['cl_pen']:.1f}%",f"{b['active']:,}",str(b['cl'])]
    for ci,v in enumerate(vals):
        c = lt.cell(ri+1,ci); c.text = v
        for p in c.text_frame.paragraphs: p.font.size = Pt(7); p.font.color.rgb = CHARCOAL

# Right — defend group
txt(sl2, Inches(6.9), Inches(1.95), Inches(6), Inches(0.25),
    f'DEFEND / DEEPEN ({len(above)} branches at or above median)', sz=10, bold=True, clr=GREEN)

rtbl = sl2.shapes.add_table(min(len(above)+1,17), 5, Inches(6.9), Inches(2.25), Inches(6.0), Inches(4.5))
rt = rtbl.table
rt.columns[0].width = Inches(2.3); rt.columns[1].width = Inches(0.9); rt.columns[2].width = Inches(0.8)
rt.columns[3].width = Inches(1.0); rt.columns[4].width = Inches(1.0)

for ci,h in enumerate(['Branch','City','CL Pen%','Active PI','CL Count']):
    c = rt.cell(0,ci); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = GREEN
    for p in c.text_frame.paragraphs: p.font.size = Pt(7); p.font.bold = True; p.font.color.rgb = WHITE

for ri,b in enumerate(sorted(above, key=lambda x: x['cl_pen'], reverse=True)[:16]):
    vals = [b['name'],b['city'],f"{b['cl_pen']:.1f}%",f"{b['active']:,}",str(b['cl'])]
    for ci,v in enumerate(vals):
        c = rt.cell(ri+1,ci); c.text = v
        for p in c.text_frame.paragraphs: p.font.size = Pt(7); p.font.color.rgb = CHARCOAL

notes2 = sl2.notes_slide
notes2.notes_text_frame.text = """SLIDE 2 — CONSUMER LOAN GROWTH

PRIORITY ACTIONS:
1. Prizren (3.2%): 11,562 active clients but only 368 consumer loans. Lowest main Filiali. Needs dedicated CL campaign.
2. Besiane/Podujeve (3.6%): 10,497 active, only 380 CL. Has 9 FTE. Capacity exists.
3. Decan (3.5%): 7,528 active, only 265 CL. High M-Banking (51.8%) = can target digitally.
4. Viti (3.5%): 5,141 active, only 182 CL. Small team (5 FTE) may be capacity-constrained.
5. Mitrovice (3.9%): 8,497 active, only 333 CL. Main Filiali with 9 FTE — underperforming.

STRONG PERFORMERS TO PROTECT:
- Bregu i Diellit (16.7%): Small branch but exceptional CL penetration
- Hani i Elezit (10.1%): Consistently strong across all metrics
- Fushe Kosove (9.8%): Strong CL despite TDA dependency issue

NOTE: CL penetration = 2025 new production / active base. This is flow intensity, not portfolio share.
"""

# ============================================================
# SLIDE 3 — FUNDING MIX & TDA STRATEGY
# ============================================================
sl3 = prs.slides.add_slide(blank)
header_bar(sl3, 'DEPOSIT FUNDING MIX \u2014 COST STRUCTURE BY BRANCH', 'CA (low-cost) vs TDA (high-cost) vs Savings | Where to moderate TDA')

exec_bar(sl3, Inches(1.1), 'SO WHAT:',
    'Three branches exceed 50% TDA share: Prizren 2 (56.8%), Gjakove (54.4%), Fushe Kosove (50.0%). '
    'At the other end, Shterpce (73.9% CA), Lipjan (70.6%), Skenderaj (69.1%) demonstrate that low-cost funding dominance is achievable. '
    'No rate data is available \u2014 cost logic is structural (TDA > SA > CA). Management should set CA/salary acquisition targets '
    'for TDA-heavy branches and stop incentivizing TDA growth where it harms economics.')

# TDA-heavy branches (left)
txt(sl3, Inches(0.4), Inches(1.95), Inches(6), Inches(0.25),
    'TDA-HEAVY BRANCHES \u2014 MODERATE TERM DEPOSIT GROWTH', sz=10, bold=True, clr=RED)

tda_high = sorted(BB, key=lambda x: x['tda_pct'], reverse=True)[:10]
ttbl = sl3.shapes.add_table(len(tda_high)+1, 5, Inches(0.4), Inches(2.25), Inches(6.0), Inches(3.5))
tt = ttbl.table
tt.columns[0].width = Inches(2.3); tt.columns[1].width = Inches(0.9); tt.columns[2].width = Inches(0.8)
tt.columns[3].width = Inches(0.8); tt.columns[4].width = Inches(1.2)

for ci,h in enumerate(['Branch','TDA%','CA%','SA%','Action']):
    c = tt.cell(0,ci); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs: p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = WHITE

for ri,b in enumerate(tda_high):
    action = 'MODERATE TDA' if b['tda_pct'] > 50 else ('MONITOR' if b['tda_pct'] > 42 else 'BALANCED')
    aclr = RED if b['tda_pct'] > 50 else (AMBER if b['tda_pct'] > 42 else GREEN)
    vals = [b['name'],f"{b['tda_pct']:.1f}%",f"{b['ca_pct']:.1f}%",f"{b['sa_pct']:.1f}%",action]
    for ci,v in enumerate(vals):
        c = tt.cell(ri+1,ci); c.text = v
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.color.rgb = aclr if ci==4 else (RED if ci==1 and b['tda_pct']>50 else CHARCOAL)
            if ci==4: p.font.bold = True

# Low-cost leaders (right)
txt(sl3, Inches(6.8), Inches(1.95), Inches(6), Inches(0.25),
    'LOW-COST FUNDING LEADERS \u2014 PROTECT & REPLICATE', sz=10, bold=True, clr=GREEN)

ca_high = sorted(BB, key=lambda x: x['ca_pct'], reverse=True)[:10]
ctbl = sl3.shapes.add_table(len(ca_high)+1, 5, Inches(6.8), Inches(2.25), Inches(6.0), Inches(3.5))
ct = ctbl.table
ct.columns[0].width = Inches(2.3); ct.columns[1].width = Inches(0.9); ct.columns[2].width = Inches(0.8)
ct.columns[3].width = Inches(0.8); ct.columns[4].width = Inches(1.2)

for ci,h in enumerate(['Branch','CA%','TDA%','SA%','Model']):
    c = ct.cell(0,ci); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs: p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = WHITE

for ri,b in enumerate(ca_high):
    model = 'BENCHMARK' if b['ca_pct'] > 65 else 'STRONG'
    vals = [b['name'],f"{b['ca_pct']:.1f}%",f"{b['tda_pct']:.1f}%",f"{b['sa_pct']:.1f}%",model]
    for ci,v in enumerate(vals):
        c = ct.cell(ri+1,ci); c.text = v
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(7); p.font.color.rgb = GREEN if ci==4 else CHARCOAL

# Warning box
shp = box(sl3, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.65), fill=RGBColor(0xFF,0xFB,0xEB), brd=AMBER)
tf = shp.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'CAUTION:  '; r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = AMBER
r2 = p.add_run(); r2.text = 'Do not reduce TDA blindly. If branch or bank liquidity requires term deposits, maintain them but flag the cost impact. '\
    'The recommendation is to STOP INCENTIVIZING additional TDA growth in branches where CA/salary alternatives are realistic. '\
    'No rate data available \u2014 all funding-cost logic is directional (TDA is structurally more expensive than CA).'; r2.font.size = Pt(8); r2.font.color.rgb = CHARCOAL

notes3 = sl3.notes_slide
notes3.notes_text_frame.text = """SLIDE 3 — FUNDING MIX

TDA MODERATION TARGETS:
- Prizren 2: 56.8% TDA — most TDA-dependent branch. Shift to CA/salary targeting.
- Gjakove: 54.4% TDA — large Filiali, 10 FTE. Should be achievable with commercial redirect.
- Fushe Kosove: 50.0% TDA — despite strong consumer lending (9.8%), funding mix is inefficient.

LOW-COST FUNDING BENCHMARKS:
- Shterpce: 73.9% CA — best in network
- Lipjan: 70.6% CA — model for ekspoziture
- Skenderaj: 69.1% CA
- Kacanik: 68.8% CA
These branches prove that >65% CA share is achievable in Kosovo's market.

DATA LIMITATION: No interest rate data exists in any source file. TDA > SA > CA cost hierarchy
is structural knowledge, not rate-calculated. If management wants specific NIM impact analysis,
rate data must be provided.
"""

# ============================================================
# SLIDE 4 — POS & MERCHANT ACQUIRING
# ============================================================
sl4 = prs.slides.add_slide(blank)
header_bar(sl4, 'POS & MERCHANT ACQUIRING \u2014 THE UNTAPPED COMMERCIAL LEVER', 'Terminal density, business coverage, deployment priorities')

exec_bar(sl4, Inches(1.1), 'SO WHAT:',
    'NLB Prishtina has 3,819 POS terminals across ~148,700 businesses \u2014 a 2.6% density. Five branches have ZERO POS terminals. '
    '1001 Prishtine serves a 21,396-business market with no merchant acquiring whatsoever. '
    'POS is not just an infrastructure metric \u2014 it is a commercial product that drives business-client acquisition, transaction revenue, and cashless migration.')

# KPIs
kpi_box(sl4, Inches(0.4), Inches(2.0), 'NETWORK POS TERMINALS', '3,819', '2.6% density across ~148,700 businesses', TEAL)
kpi_box(sl4, Inches(4.5), Inches(2.0), 'BRANCHES WITH ZERO POS', '5', '1001 Prishtine, Prizren 2, Llapit, B.Diellit, Mit.2', RED)
kpi_box(sl4, Inches(8.6), Inches(2.0), 'TOP POS DENSITY', '10.1%', 'Filiali Peje (421 POS / 4,172 businesses)', GREEN)

# POS priority table
txt(sl4, Inches(0.4), Inches(3.1), Inches(12), Inches(0.25),
    'POS DEPLOYMENT PRIORITIES', sz=11, bold=True, clr=NAVY)

pos_sorted = sorted(BB, key=lambda x: x['pos_den'])
pos_rows = []
for b in pos_sorted[:15]:
    action = 'DEPLOY IMMEDIATELY' if b['pos'] == 0 else ('ACCELERATE' if b['pos_den'] < 2 else 'GROW')
    pos_rows.append([b['name'],b['city'],str(b['pos']),f"{b['biz']:,}",f"{b['pos_den']:.1f}%",action])

ptbl = sl4.shapes.add_table(len(pos_rows)+1, 6, Inches(0.4), Inches(3.4), Inches(12.5), Inches(3.5))
pt = ptbl.table
pt.columns[0].width = Inches(2.8); pt.columns[1].width = Inches(1.5); pt.columns[2].width = Inches(1.0)
pt.columns[3].width = Inches(1.5); pt.columns[4].width = Inches(1.5); pt.columns[5].width = Inches(4.2)

for ci,h in enumerate(['Branch','City','POS Count','Businesses','POS/100 Biz','Action']):
    c = pt.cell(0,ci); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs: p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = WHITE

for ri,row in enumerate(pos_rows):
    for ci,v in enumerate(row):
        c = pt.cell(ri+1,ci); c.text = v
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(7)
            if ci==5:
                p.font.bold = True
                p.font.color.rgb = RED if 'IMMEDIATELY' in v else (AMBER if 'ACCELERATE' in v else GREEN)
            elif ci==4 and v=='0.0%':
                p.font.color.rgb = RED; p.font.bold = True
            else:
                p.font.color.rgb = CHARCOAL

notes4 = sl4.notes_slide
notes4.notes_text_frame.text = """SLIDE 4 — POS & MERCHANT ACQUIRING

ZERO-POS BRANCHES (DEPLOY IMMEDIATELY):
1. 1001 Prishtine: 21,396 businesses, 18 FTE, ZERO POS. Largest single gap.
2. Prizren 2 (100504): 7,119 businesses, 6 FTE, ZERO POS.
3. Prishtine-Llapit (000102): Shares Prishtine market, 6 FTE, ZERO POS.
4. Bregu i Diellit (000110): Shares Prishtine market, 4 FTE, ZERO POS.
5. Mitrovice 2 (100704): 2,529 businesses, 5 FTE, ZERO POS.

TOP POS DENSITY (DEFEND & REPLICATE):
- Peje: 10.1% (421/4,172) — best in network
- Hani i Elezit: 9.6% (27/282)
- Prishtine 0001: 8.0% (1,705/21,396) — shows Prishtine CAN do merchant acquiring
- Fushe Kosove: 7.4% (201/2,711)

POS GROWTH STRATEGY:
- POS is a business-client acquisition tool, not just payments infrastructure
- Each POS terminal creates a recurring revenue relationship
- POS density correlates with business client retention and cashless transaction share
- Recommend minimum 3% density target for all branches within 12 months
"""

# Save
fname = 'Module2_Product_Mix.pptx'
try:
    prs.save(fname)
except PermissionError:
    fname = 'Module2_Product_Mix_v2.pptx'
    prs.save(fname)

print(f"Saved: {fname}")
print(f"Size: {os.path.getsize(fname):,} bytes")
print(f"Slides: {len(prs.slides)}")
