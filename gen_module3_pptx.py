#!/usr/bin/env python3
"""
NLB Prishtina — Module 3: Staff Efficiency & Productivity
5-Slide Executive PPTX | McKinsey-style
Linked to Module 1 (Market Share) and Module 2 (Product Mix)
"""
import sys, json, os, statistics
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# DESIGN TOKENS
# ============================================================
NAVY = RGBColor(0x00,0x2B,0x49)
STEEL = RGBColor(0x4A,0x6F,0xA5)
CHARCOAL = RGBColor(0x1A,0x1A,0x2E)
WHITE = RGBColor(0xFF,0xFF,0xFF)
GREEN = RGBColor(0x05,0x96,0x69)
RED = RGBColor(0xDC,0x26,0x26)
AMBER = RGBColor(0xD9,0x77,0x06)
TEAL = RGBColor(0x0D,0x94,0x88)
GRAY = RGBColor(0x64,0x74,0x8B)
BORDER = RGBColor(0xE2,0xE8,0xF0)
LIGHT = RGBColor(0xF0,0xF4,0xF8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============================================================
# LOAD DATA
# ============================================================
with open('branch_data_unified.json', encoding='utf-8') as f:
    data = json.load(f)
branches = data['branches']

POS = {'0001':1705,'1006':421,'100104':201,'1005':156,'100106':126,'1007':107,
       '1002':115,'1003':107,'100603':71,'100702':69,'000103':78,
       '100201':50,'100602':62,'100701':47,'1009':45,'1004':60,'100503':56,
       '100501':43,'100604':54,'100302':42,'100502':33,'100301':27,'100202':28,
       '100105':11,'100303':9,'100108':6,'000102':0,'000110':0,'1001':0,'100504':0,'100704':0}

def s(v, d=0):
    if v is None: return d
    try: return float(v)
    except: return d

# Build benchmarked branch list
BB = []
for ou, b in branches.items():
    if not b.get('is_benchmarked'):
        continue
    emp = b['employees'] or 1
    active = s(b.get('active_clients_pi'))
    total_txn = s(b.get('total_transactions'))
    cash_txn = s(b.get('total_cash_transactions'))
    pi_ob = s(b.get('pi_onbalance_segment'))
    micro_ob = s(b.get('micro_onbalance_segment'))
    cib_ob = s(b.get('cib_onbalance_segment'))
    total_ob = s(b.get('total_onbalance'))
    cons = s(b.get('consumer_loan_no'))
    micro_loans = s(b.get('micro_loan_count'))
    cib_loans = s(b.get('cib_loan_count'))
    pi_loans = s(b.get('total_pi_loans'))
    new_acc = s(b.get('new_accounts'))
    deposits = s(b.get('total_deposits'))
    daily_svc = s(b.get('daily_services_completed'))
    xsell = s(b.get('cross_sell_intensity'))
    cash_r = cash_txn / total_txn if total_txn > 0 else 0
    pi_pen = s(b.get('pop_penetration')) * 100

    total_loans = pi_loans + micro_loans + cib_loans
    pos = POS.get(ou, 0)

    # M1 group
    if pi_pen >= 18.0:
        m1 = 'Defend'
    elif pi_pen >= 13.1:
        m1 = 'Acceleration'
    else:
        m1 = 'Recovery'

    BB.append({
        'ou': ou, 'name': b['branch_name'].replace('Filiali ','').replace('Ekspozitura ',''),
        'city': b.get('city',''), 'type': b.get('branch_type',''),
        'fte': emp, 'actPI': active, 'inactPI': s(b.get('non_active_clients_pi')),
        'totalTxn': total_txn, 'cashTxn': cash_txn, 'cashR': cash_r,
        'piLoans': pi_loans, 'consLoans': cons, 'microLoans': micro_loans, 'cibLoans': cib_loans,
        'totalLoans': total_loans,
        'piOB': pi_ob/1e6, 'microOB': micro_ob/1e6, 'cibOB': cib_ob/1e6, 'totalOB': total_ob/1e6,
        'newAcc': new_acc, 'deposits': deposits/1e6, 'dailySvc': daily_svc,
        'xsell': xsell, 'piPen': pi_pen, 'pos': pos,
        # Per FTE
        'piClientsFTE': active/emp, 'txnFTE': total_txn/emp,
        'consLoansFTE': cons/emp, 'loansFTE': total_loans/emp,
        'newAccFTE': new_acc/emp, 'obFTE': total_ob/emp/1e6,
        'piOBfte': pi_ob/emp/1e6, 'microOBfte': micro_ob/emp/1e6,
        'cibOBfte': cib_ob/emp/1e6, 'depFTE': deposits/emp/1e6,
        'cashFTE': cash_txn/emp, 'posFTE': pos/emp,
        'microLoansFTE': micro_loans/emp,
        'm1': m1,
    })

# Compute medians
def med(vals):
    v = [x for x in vals if x > 0]
    return statistics.median(v) if v else 0

MEDS = {
    'piClientsFTE': med([b['piClientsFTE'] for b in BB]),
    'consLoansFTE': med([b['consLoansFTE'] for b in BB]),
    'loansFTE': med([b['loansFTE'] for b in BB]),
    'newAccFTE': med([b['newAccFTE'] for b in BB]),
    'obFTE': med([b['obFTE'] for b in BB]),
    'cashR': med([b['cashR'] for b in BB]),
    'dailySvc': med([b['dailySvc'] for b in BB if b['dailySvc'] > 0]),
    'txnFTE': med([b['txnFTE'] for b in BB]),
    'xsell': med([b['xsell'] for b in BB]),
    'microOBfte': med([b['microOBfte'] for b in BB if b['microOBfte'] > 0]),
    'depFTE': med([b['depFTE'] for b in BB]),
}

# Compute composite score and targets
for b in BB:
    norm = lambda v, m: min(v/m, 2) if m > 0 else 0
    b['prodScore'] = (norm(b['piClientsFTE'], MEDS['piClientsFTE']) * 0.20 +
                      norm(b['consLoansFTE'], MEDS['consLoansFTE']) * 0.20 +
                      norm(b['newAccFTE'], MEDS['newAccFTE']) * 0.15 +
                      norm(b['obFTE'], MEDS['obFTE']) * 0.15 +
                      norm(b['microLoansFTE'], MEDS.get('microLoansFTE', 1) or 1) * 0.10 +
                      (1 - b['cashR']) * 0.10 +
                      norm(b['xsell'], MEDS['xsell']) * 0.10)

    # Targets
    b['tgtConsFTE'] = max(round(MEDS['consLoansFTE'] * 1.05), round(b['consLoansFTE'])) if b['consLoansFTE'] >= MEDS['consLoansFTE'] else round(MEDS['consLoansFTE'] * 1.05)
    if b['m1'] == 'Recovery':
        b['tgtConsFTE'] = max(b['tgtConsFTE'], round(MEDS['consLoansFTE'] * 1.15))

    b['tgtCashR'] = 0.60 if b['cashR'] > 0.75 else (b['cashR'] - 0.10 if b['cashR'] > 0.60 else b['cashR'])
    b['cashReduction'] = max(0, b['cashR'] - b['tgtCashR'])

    # Role-mix
    sales_idx = (b['consLoansFTE'] + b['newAccFTE']) / (MEDS['consLoansFTE'] + MEDS['newAccFTE']) if (MEDS['consLoansFTE'] + MEDS['newAccFTE']) > 0 else 0
    if b['cashR'] > 0.75 and sales_idx < 0.9:
        b['roleMix'] = 'Teller-Heavy / Sales-Light'
    elif b['cashR'] > 0.70 and sales_idx >= 0.9:
        b['roleMix'] = 'Balanced but Cash-Burdened'
    elif sales_idx > 1.2:
        b['roleMix'] = 'Sales-Strong'
    elif b['microOB'] == 0 and b['cibOB'] == 0:
        b['roleMix'] = 'PI-Only / No Business'
    else:
        b['roleMix'] = 'Adequate'

    # Action
    if b['cashR'] > 0.80 and b['piClientsFTE'] > MEDS['piClientsFTE']:
        b['action'] = 'Digital Migration'
    elif b['piClientsFTE'] > MEDS['piClientsFTE']*1.15 and b['consLoansFTE'] > MEDS['consLoansFTE']*1.15:
        b['action'] = 'Hire / Reinforce'
    elif b['cashR'] > 0.75 and b['consLoansFTE'] < MEDS['consLoansFTE']*0.8:
        b['action'] = 'Role Redesign + Digital'
    elif b['consLoansFTE'] < MEDS['consLoansFTE']*0.8 or b['newAccFTE'] < MEDS['newAccFTE']*0.8:
        b['action'] = 'Sales Discipline'
    elif b['obFTE'] > MEDS['obFTE']*1.3:
        b['action'] = 'Defend Stock + Push Flow'
    else:
        b['action'] = 'Maintain + Optimize'

BB.sort(key=lambda x: x['prodScore'], reverse=True)

# ============================================================
# PPTX BUILDER
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_shape(slide, left, top, w, h, fill=None, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.line.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.75)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    return shp

def add_textbox(slide, left, top, w, h, text='', size=10, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = 'Inter'
    return tb

def kpi_box(slide, left, top, w, h, label, value, context='', accent=NAVY):
    shp = add_shape(slide, left, top, w, h, fill=WHITE, border=BORDER)
    # Left accent bar
    bar = add_shape(slide, left, top, Inches(0.06), h, fill=accent)
    # Label
    add_textbox(slide, left+Inches(0.15), top+Inches(0.08), w-Inches(0.2), Inches(0.25),
                label, size=8, bold=True, color=GRAY)
    # Value
    add_textbox(slide, left+Inches(0.15), top+Inches(0.30), w-Inches(0.2), Inches(0.35),
                value, size=22, bold=True, color=NAVY)
    # Context
    if context:
        add_textbox(slide, left+Inches(0.15), top+Inches(0.65), w-Inches(0.2), Inches(0.25),
                    context, size=7, color=GRAY)

def slide_header(slide, title, subtitle=''):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), fill=NAVY)
    add_textbox(slide, Inches(0.5), Inches(0.12), Inches(8), Inches(0.4),
                title, size=18, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.50), Inches(10), Inches(0.3),
                    subtitle, size=9, color=RGBColor(0xA0,0xB0,0xC0))
    # Module badge
    add_textbox(slide, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.3),
                'MODULE 3 | Staff Efficiency', size=8, bold=True, color=RGBColor(0x80,0x99,0xB0), align=PP_ALIGN.RIGHT)

def so_what(slide, left, top, w, text):
    shp = add_shape(slide, left, top, w, Inches(0.45), fill=RGBColor(0xEF,0xF6,0xFF), border=STEEL)
    tb = slide.shapes.add_textbox(left+Inches(0.12), top+Inches(0.05), w-Inches(0.24), Inches(0.35))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    r1 = p.add_run()
    r1.text = 'SO WHAT: '
    r1.font.size = Pt(8)
    r1.font.bold = True
    r1.font.color.rgb = STEEL
    r1.font.name = 'Inter'
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(8)
    r2.font.color.rgb = CHARCOAL
    r2.font.name = 'Inter'

def add_table(slide, left, top, w, h, headers, rows, col_widths=None):
    tbl_shape = slide.shapes.add_table(len(rows)+1, len(headers), left, top, w, h)
    tbl = tbl_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    # Header row
    for i, hdr in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(7)
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.name = 'Inter'
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(7)
                    r.font.color.rgb = CHARCOAL
                    r.font.name = 'Inter'
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return tbl_shape

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

# ============================================================
# SLIDE 1: ACTUAL PRODUCTIVITY OVERVIEW
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide1, 'Staff Productivity — Actual Performance', 'Part 1: What staff are producing today | 31 benchmarked branches | 265 FTE')

# KPI boxes
kw = Inches(2.35)
kh = Inches(0.9)
y0 = Inches(1.05)
kpi_box(slide1, Inches(0.3), y0, kw, kh, 'PI CLIENTS / FTE', f'{MEDS["piClientsFTE"]:.0f}', 'Median across 31 branches', NAVY)
kpi_box(slide1, Inches(2.8), y0, kw, kh, 'CONSUMER LOANS / FTE', f'{MEDS["consLoansFTE"]:.1f}', 'Median — 2025 flow', GREEN)
kpi_box(slide1, Inches(5.3), y0, kw, kh, 'CASH TRANSACTION SHARE', f'{MEDS["cashR"]*100:.0f}%', 'Median — teller burden proxy', RED)
kpi_box(slide1, Inches(7.8), y0, kw, kh, 'DAILY SERVICES / FTE', f'{MEDS["dailySvc"]:.1f}', 'Median (28 branches)', AMBER)
kpi_box(slide1, Inches(10.3), y0, kw, kh, 'ONBALANCE / FTE', f'EUR {MEDS["obFTE"]:.1f}M', 'Median stock depth', STEEL)

so_what(slide1, Inches(0.3), Inches(2.1), Inches(12.5),
        '73% of branch transactions are cash-based. This is the single largest drag on staff productivity. '
        '12 branches exceed 75% cash share. Released teller capacity should be redirected toward consumer lending, '
        'client activation, and advisory. The target is not more manual transactions, but fewer.')

# Productivity ranking table — top 15 + bottom 5
top15 = BB[:15]
bot5 = BB[-5:]

headers = ['Branch', 'FTE', 'PI Cl/FTE', 'Cons/FTE', 'NewAcc/FTE', 'Loans/FTE', 'OB/FTE M', 'Cash%', 'X-Sell', 'Daily Svc', 'Score']
rows_data = []
for b in top15:
    rows_data.append([
        b['name'], str(b['fte']),
        f'{b["piClientsFTE"]:.0f}', f'{b["consLoansFTE"]:.1f}',
        f'{b["newAccFTE"]:.0f}', f'{b["loansFTE"]:.0f}',
        f'{b["obFTE"]:.1f}', f'{b["cashR"]*100:.0f}%',
        f'{b["xsell"]:.3f}', f'{b["dailySvc"]:.1f}' if b['dailySvc'] > 0 else '—',
        f'{b["prodScore"]:.2f}'
    ])
# Separator
rows_data.append(['--- Bottom 5 ---', '', '', '', '', '', '', '', '', '', ''])
for b in bot5:
    rows_data.append([
        b['name'], str(b['fte']),
        f'{b["piClientsFTE"]:.0f}', f'{b["consLoansFTE"]:.1f}',
        f'{b["newAccFTE"]:.0f}', f'{b["loansFTE"]:.0f}',
        f'{b["obFTE"]:.1f}', f'{b["cashR"]*100:.0f}%',
        f'{b["xsell"]:.3f}', f'{b["dailySvc"]:.1f}' if b['dailySvc'] > 0 else '—',
        f'{b["prodScore"]:.2f}'
    ])

add_table(slide1, Inches(0.3), Inches(2.7), Inches(12.7), Inches(4.5), headers, rows_data,
          col_widths=[1.8, 0.5, 0.9, 0.9, 0.9, 0.9, 0.9, 0.7, 0.7, 0.8, 0.7])

# Link badges
add_textbox(slide1, Inches(0.3), Inches(7.1), Inches(6), Inches(0.3),
            'Linked: M1 Market Share groups | M2 Product Mix priorities', size=7, color=STEEL)

add_notes(slide1, """SLIDE 1 — ACTUAL PRODUCTIVITY OVERVIEW

TALKING POINTS:
1. Network median: 914 PI clients per FTE, 48 consumer loans per FTE, EUR 3.1M OnBalance per FTE
2. Critical finding: 73% median cash transaction share — the dominant share of staff time is consumed by low-value manual cash handling
3. 12 branches exceed 75% cash ratio, with Mitrovice 2 (93%), Skenderaj (84%), Obiliq (84%), Decan (80%) as worst cases
4. Top performers: Bregu i Diellit leads on flow intensity (102 cons loans/FTE), Prishtine-Llapit on new accounts (220/FTE)
5. Bottom performers: Ferizaj (19 FTE) scores lowest among main branches — high headcount but below-median consumer lending and high cash burden
6. The productivity gap is not primarily about headcount — it's about what staff are spending their time on
7. Daily services range from 10.0 (Kacanik) to 43.3 (Prishtine) — a 4x spread that signals significant utilization differences

MANAGEMENT MESSAGE:
The network's productivity problem is primarily a workload-composition problem, not a staffing-level problem. Staff are spending most of their time on cash transactions instead of commercial activity. Before adding headcount anywhere, the first priority is migrating low-value transactions to digital channels to release capacity for selling.""")

# ============================================================
# SLIDE 2: CASH BURDEN & TRANSACTION ANALYSIS
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide2, 'Transaction Burden — Cash vs Digital', 'The single largest drag on staff productivity is manual cash handling')

# KPIs
y0 = Inches(1.05)
cash_heavy = [b for b in BB if b['cashR'] > 0.75]
low_cash = [b for b in BB if b['cashR'] < 0.50]
avg_cash_fte = statistics.mean([b['cashFTE'] for b in BB])

kpi_box(slide2, Inches(0.3), y0, Inches(2.9), kh, 'BRANCHES > 75% CASH', f'{len(cash_heavy)} of 31', 'Teller-dominated workload', RED)
kpi_box(slide2, Inches(3.4), y0, Inches(2.9), kh, 'BRANCHES < 50% CASH', f'{len(low_cash)} of 31', 'Healthier digital balance', GREEN)
kpi_box(slide2, Inches(6.5), y0, Inches(2.9), kh, 'AVG CASH TXN / FTE', f'{avg_cash_fte:.0f}', 'Annual cash transactions per employee', AMBER)
kpi_box(slide2, Inches(9.6), y0, Inches(2.9), kh, 'POTENTIAL FTE FREED', f'~15-20', 'If cash >75% reduced to 60%', TEAL)

so_what(slide2, Inches(0.3), Inches(2.1), Inches(12.5),
        'The target is not more manual transactions but more cashless migration. '
        'If the 12 branches above 75% cash ratio reduced to 60%, approximately 15-20 FTE-equivalents of capacity '
        'would be released for consumer lending, PI acquisition, and inactive-client activation.')

# Cash burden table — branches sorted by cash ratio descending
cash_sorted = sorted(BB, key=lambda x: x['cashR'], reverse=True)[:20]
headers2 = ['Branch', 'FTE', 'Total Txn', 'Cash Txn', 'Cash %', 'Cash/FTE', 'Non-Cash Txn',
            'Cons Loans/FTE', 'M1 Group', 'Diagnostic']
rows2 = []
for b in cash_sorted:
    non_cash = b['totalTxn'] - b['cashTxn']
    if b['cashR'] > 0.80:
        diag = 'Critical — mandatory digital push'
    elif b['cashR'] > 0.75:
        diag = 'High — accelerate cashless shift'
    elif b['cashR'] > 0.65:
        diag = 'Moderate — continue migration'
    else:
        diag = 'Acceptable'
    rows2.append([
        b['name'], str(b['fte']),
        f'{b["totalTxn"]:,.0f}', f'{b["cashTxn"]:,.0f}',
        f'{b["cashR"]*100:.0f}%', f'{b["cashFTE"]:.0f}',
        f'{non_cash:,.0f}',
        f'{b["consLoansFTE"]:.1f}', b['m1'],
        diag
    ])

add_table(slide2, Inches(0.3), Inches(2.7), Inches(12.7), Inches(4.5), headers2, rows2,
          col_widths=[1.6, 0.5, 0.9, 0.9, 0.7, 0.7, 0.9, 0.9, 1.0, 2.2])

add_notes(slide2, """SLIDE 2 — CASH BURDEN & TRANSACTION ANALYSIS

TALKING POINTS:
1. Mitrovice 2 has 93% cash ratio — nearly all branch activity is manual cash handling. With only 5 FTE, the branch has no capacity for commercial activity.
2. Skenderaj (84%), Obiliq (84%), Decan (80%), Suhareke (80%), Mitrovice (81%) — all above 80% threshold
3. The 12 branches above 75% collectively process ~85,000 cash transactions per year. Migrating even 20% to digital channels would release significant staff time.
4. Contrast with Kline (37% cash) and Prishtine (28% cash) — these branches have a healthier digital balance and higher commercial output per FTE
5. The FTE freed estimate (~15-20) is directional: (cash reduction * branch txn volume) / median txn per FTE
6. Released capacity should be redirected toward: consumer lending push, new PI account acquisition, inactive client activation, POS/merchant development

MANAGEMENT MESSAGE:
Cash transaction intensity is the primary bottleneck. The issue is not that branches need more tellers — the issue is that teller work consumes too much of the branch. The path to higher productivity runs through digital migration, not additional headcount. Every 10% reduction in cash share at a branch with 8 FTE releases approximately 0.8 FTE-equivalent for commercial activity.""")

# ============================================================
# SLIDE 3: SEGMENT PRODUCTIVITY — PI, MICRO, CIB
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide3, 'Segment Productivity per FTE', 'PI lending, Micro portfolio depth, and business segment coverage by branch')

y0 = Inches(1.05)
kpi_box(slide3, Inches(0.3), y0, Inches(2.0), kh, 'PI OB / FTE', f'EUR {med([b["piOBfte"] for b in BB]):.1f}M', 'Median PI portfolio depth', NAVY)
kpi_box(slide3, Inches(2.5), y0, Inches(2.0), kh, 'MICRO OB / FTE', f'EUR {MEDS["microOBfte"]:.2f}M', 'Median micro depth', TEAL)
kpi_box(slide3, Inches(4.7), y0, Inches(2.0), kh, 'CONS LOANS / FTE', f'{MEDS["consLoansFTE"]:.1f}', 'Median flow output', GREEN)
kpi_box(slide3, Inches(6.9), y0, Inches(2.0), kh, 'NEW ACCOUNTS / FTE', f'{MEDS["newAccFTE"]:.0f}', 'Median acquisition rate', AMBER)
kpi_box(slide3, Inches(9.1), y0, Inches(2.0), kh, 'CROSS-SELL', f'{MEDS["xsell"]:.3f}', 'Median intensity', STEEL)
kpi_box(slide3, Inches(11.3), y0, Inches(2.0), kh, 'DEPOSITS / FTE', f'EUR {MEDS["depFTE"]:.1f}M', 'Median deposit depth', GRAY)

so_what(slide3, Inches(0.3), Inches(2.1), Inches(12.5),
        'PI lending productivity varies 3x across branches. Branches like Ferizaj (19 FTE) underperform smaller branches on consumer loans/FTE. '
        'Micro portfolio depth is concentrated in a few branches — 5 branches have zero or near-zero micro presence. '
        'CIB is dominated by Prishtine main branches. This is a role-allocation problem, not just a staffing problem.')

# Segment table — all 31 branches
headers3 = ['Branch', 'FTE', 'Type', 'PI OB/FTE', 'Cons/FTE', 'NewAcc/FTE',
            'Micro OB/FTE', 'Micro Ln/FTE', 'CIB OB/FTE', 'X-Sell', 'Role Mix']
rows3 = []
for b in BB[:20]:  # Top 20 by score
    rows3.append([
        b['name'], str(b['fte']), b['type'][:8],
        f'{b["piOBfte"]:.2f}', f'{b["consLoansFTE"]:.1f}', f'{b["newAccFTE"]:.0f}',
        f'{b["microOBfte"]:.2f}' if b['microOBfte'] > 0 else '—',
        f'{b["microLoansFTE"]:.1f}' if b['microLoansFTE'] > 0 else '—',
        f'{b["cibOBfte"]:.2f}' if b['cibOBfte'] > 0 else '—',
        f'{b["xsell"]:.3f}', b['roleMix'][:20]
    ])

add_table(slide3, Inches(0.3), Inches(2.7), Inches(12.7), Inches(4.5), headers3, rows3,
          col_widths=[1.5, 0.4, 0.7, 0.8, 0.8, 0.8, 0.8, 0.8, 0.8, 0.7, 1.8])

add_textbox(slide3, Inches(0.3), Inches(7.1), Inches(6), Inches(0.3),
            'Note: All "per employee" metrics use total branch FTE — no role-level headcount breakdown available', size=7, color=GRAY)

add_notes(slide3, """SLIDE 3 — SEGMENT PRODUCTIVITY PER FTE

TALKING POINTS:
1. PI portfolio depth: ranges from EUR 0.86M/FTE (Gracanice) to EUR 5.52M/FTE (1001 Prishtine). Top-heavy concentration.
2. Consumer loans per FTE: Bregu i Diellit leads (102/FTE) despite having only 4 FTE. Ferizaj (41/FTE) underperforms with 19 FTE.
3. Micro presence is highly uneven: Kacanik leads on micro loans/FTE (18.0), Lipjan on micro portfolio/FTE (EUR 0.78M). 5 branches have zero micro activity.
4. CIB is concentrated: Prishtine 0001 (EUR 8.98M/FTE) and 1001 Prishtine (EUR 5.78M/FTE) dominate. Most ekspoziture have negligible CIB.
5. Cross-sell intensity median is only 0.26 — meaning most branches sell fewer than 0.3 products per active client per year. Bregu i Diellit (0.70) is the outlier.
6. Role-mix assessment: 8 branches are "Teller-Heavy / Sales-Light", 5 are "PI-Only / No Business".

MANAGEMENT MESSAGE:
Segment productivity reveals that the issue is not just total output but what staff are producing. Several large branches are underperforming smaller ones on commercial metrics. The network needs role-mix redesign — fewer teller-hours, more advisory and sales time — especially in Recovery branches where market-share gaps demand higher acquisition rates.""")

# ============================================================
# SLIDE 4: ACTUAL vs TARGET PRODUCTIVITY
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide4, 'Target Productivity — Closing the Gap', 'Part 2: What staff should produce next | Targets linked to M1 Market Share & M2 Product Mix')

y0 = Inches(1.05)
recovery = [b for b in BB if b['m1'] == 'Recovery']
acceleration = [b for b in BB if b['m1'] == 'Acceleration']
defend = [b for b in BB if b['m1'] == 'Defend']

kpi_box(slide4, Inches(0.3), y0, Inches(2.9), kh, 'RECOVERY BRANCHES (M1)', f'{len(recovery)}', 'Need higher acquisition + lending targets', RED)
kpi_box(slide4, Inches(3.4), y0, Inches(2.9), kh, 'ACCELERATION BRANCHES (M1)', f'{len(acceleration)}', 'Push from median toward 18% target', AMBER)
kpi_box(slide4, Inches(6.5), y0, Inches(2.9), kh, 'DEFEND BRANCHES (M1)', f'{len(defend)}', 'Maintain depth + quality focus', GREEN)
kpi_box(slide4, Inches(9.6), y0, Inches(2.9), kh, 'ROLE REDESIGN NEEDED', f'{len([b for b in BB if "Redesign" in b["action"] or "Digital" in b["action"]])}', 'Cash-to-sales rebalancing', TEAL)

so_what(slide4, Inches(0.3), Inches(2.1), Inches(12.5),
        'Targets are derived from market-share gaps (M1), product-mix priorities (M2), and peer-group medians. '
        'Recovery branches get 15% uplift targets on consumer lending. '
        'Cash-heavy branches get reduction targets, not transaction-increase targets. '
        'The freed capacity from digital migration is the primary source of productivity improvement.')

# Target comparison table
targets_sorted = sorted(BB, key=lambda x: x['cashR'], reverse=True)[:18]
headers4 = ['Branch', 'FTE', 'M1 Grp', 'Cons/FTE Now', 'Cons/FTE Tgt', 'Gap',
            'Cash% Now', 'Cash% Tgt', 'FTE Freed*', 'Action']
rows4 = []
for b in targets_sorted:
    cons_gap = b['tgtConsFTE'] - b['consLoansFTE']
    freed = b['cashReduction'] * b['totalTxn'] / MEDS['txnFTE'] if MEDS['txnFTE'] > 0 else 0
    rows4.append([
        b['name'], str(b['fte']), b['m1'],
        f'{b["consLoansFTE"]:.1f}', f'{b["tgtConsFTE"]}',
        f'+{cons_gap:.0f}' if cons_gap > 0 else f'{cons_gap:.0f}',
        f'{b["cashR"]*100:.0f}%', f'{b["tgtCashR"]*100:.0f}%',
        f'{freed:.1f}' if freed > 0 else '—',
        b['action']
    ])

add_table(slide4, Inches(0.3), Inches(2.7), Inches(12.7), Inches(4.2), headers4, rows4,
          col_widths=[1.5, 0.5, 0.8, 0.9, 0.9, 0.6, 0.8, 0.8, 0.7, 2.2])

add_textbox(slide4, Inches(0.3), Inches(7.0), Inches(12), Inches(0.35),
            '* FTE Freed = directional estimate of capacity released if cash ratio reduced to target level. '
            'Not a headcount reduction — redeploy to commercial activity.', size=7, color=GRAY)

add_notes(slide4, """SLIDE 4 — TARGET PRODUCTIVITY

TALKING POINTS:
1. Recovery branches (PI pen < 13.1%) need consumer lending targets 15% above current median — not because current output is zero, but because market-share gaps demand acceleration
2. Acceleration branches need incremental push from median toward 18% target — steady improvement, not dramatic shift
3. Defend branches should focus on stock quality, relationship deepening, and cross-sell rather than volume growth
4. Cash reduction targets: branches above 75% → target 60%. This is achievable through POS migration, M-banking push, and branch process redesign
5. The "FTE Freed" column is the most important number on this slide — it shows the productivity dividend from digital migration
6. No branch gets a "do more manual transactions" target. The transaction target is reduction, not growth.
7. Role redesign branches need teller-to-advisor rebalancing, not additional headcount
8. Total network FTE freed from cash reduction: approximately 15-20 FTE-equivalents — enough to staff 3-4 additional sales positions without any new hiring

MANAGEMENT MESSAGE:
Target productivity must be achieved primarily through workload redesign, not headcount increase. The single highest-impact lever is cash-to-digital migration. Every branch above 75% cash share should have a 90-day cashless migration plan. The second lever is sales discipline — clear individual targets for consumer lending, new account acquisition, and inactive-client activation. Hiring should only be considered for branches that are already productive AND overloaded.""")

# ============================================================
# SLIDE 5: MANAGEMENT RECOMMENDATIONS
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide5, 'Management Actions — Branch-by-Branch', 'Explicit recommendations linked to M1 Market Share, M2 Product Mix, and Part 1 diagnostic')

y0 = Inches(1.05)
# Summary action count
hire_count = len([b for b in BB if 'Hire' in b['action']])
redesign_count = len([b for b in BB if 'Redesign' in b['action'] or 'Digital' in b['action']])
sales_count = len([b for b in BB if 'Sales' in b['action']])
maintain_count = len([b for b in BB if 'Maintain' in b['action'] or 'Defend' in b['action']])

kpi_box(slide5, Inches(0.3), y0, Inches(2.9), kh, 'HIRE / REINFORCE', f'{hire_count} branches', 'Strong productivity + overloaded', NAVY)
kpi_box(slide5, Inches(3.4), y0, Inches(2.9), kh, 'ROLE REDESIGN / DIGITAL', f'{redesign_count} branches', 'Cash-heavy + weak commercial output', RED)
kpi_box(slide5, Inches(6.5), y0, Inches(2.9), kh, 'SALES DISCIPLINE', f'{sales_count} branches', 'Below-median sales + adequate capacity', AMBER)
kpi_box(slide5, Inches(9.6), y0, Inches(2.9), kh, 'MAINTAIN / OPTIMIZE', f'{maintain_count} branches', 'Adequate performance + monitor', GREEN)

so_what(slide5, Inches(0.3), Inches(2.1), Inches(12.5),
        f'Only {hire_count} branches justify additional hiring — the rest need workload redesign, '
        f'digital migration, or stronger sales discipline. The network has sufficient total FTE; '
        f'the issue is allocation and activity composition, not raw headcount.')

# Recommendations table
headers5 = ['Branch', 'FTE', 'M1', 'Cash%', 'Cons/FTE', 'Score', 'Action', 'Recommendation']
rows5 = []
for b in sorted(BB, key=lambda x: {'Hire / Reinforce':0,'Defend Stock + Push Flow':1,
                                     'Role Redesign + Digital':2,'Digital Migration':3,
                                     'Sales Discipline':4,'Maintain + Optimize':5}.get(x['action'],5)):
    # Build recommendation text
    if 'Hire' in b['action']:
        reco = f'Strong output ({b["consLoansFTE"]:.0f} cons/FTE, {b["piClientsFTE"]:.0f} PI/FTE). Consider +1 FTE for {"PI acquisition" if b["m1"]=="Recovery" else "commercial deepening"}.'
    elif 'Digital' in b['action'] and 'Redesign' not in b['action']:
        reco = f'{b["cashR"]*100:.0f}% cash ratio unsustainable. Migrate to digital. Released capacity to consumer lending.'
    elif 'Redesign' in b['action']:
        reco = f'High cash ({b["cashR"]*100:.0f}%) + weak sales ({b["consLoansFTE"]:.0f} cons/FTE). Shift teller hours to advisory.'
    elif 'Sales' in b['action']:
        reco = f'Below-median output. Set individual targets. {"Reactivation pool: " + str(int(b["inactPI"])) + " inactive PI." if b["inactPI"]>1000 else ""}'
    elif 'Defend' in b['action']:
        reco = f'Strong stock (EUR {b["obFTE"]:.1f}M/FTE) but flow needs push. Focus on new production volume.'
    else:
        reco = f'Adequate performance. Monitor trends. {"Focus on retention." if b["m1"]=="Defend" else "Incremental improvement."}'

    rows5.append([
        b['name'], str(b['fte']), b['m1'],
        f'{b["cashR"]*100:.0f}%', f'{b["consLoansFTE"]:.1f}',
        f'{b["prodScore"]:.2f}', b['action'], reco
    ])

add_table(slide5, Inches(0.3), Inches(2.7), Inches(12.7), Inches(4.5), headers5, rows5[:21],
          col_widths=[1.3, 0.4, 0.8, 0.6, 0.7, 0.6, 1.8, 3.5])

add_textbox(slide5, Inches(0.3), Inches(7.1), Inches(12), Inches(0.3),
            'Full 31-branch detail in HTML Branch Intelligence Dashboard | Module 3', size=7, color=STEEL)

add_notes(slide5, """SLIDE 5 — MANAGEMENT RECOMMENDATIONS

TALKING POINTS:
1. Network-wide: the issue is workload composition, not total headcount. 265 FTE is adequate if cash-to-digital migration is executed.
2. Hire/Reinforce: Only justified for branches with proven productivity above median AND visible capacity strain. Do not hire into branches with low output.
3. Role Redesign + Digital Migration: The largest group. These branches have staff spending 75-93% of time on cash transactions. No amount of sales targets will work without reducing the cash burden first.
4. Sales Discipline: These branches have adequate capacity but below-median commercial output. The fix is individual performance targets, weekly tracking, and management accountability — not structural change.
5. Maintain: These are performing adequately. Continue monitoring, set incremental targets.
6. Priority sequence: (a) Deploy cashless migration in 12 cash-heavy branches immediately, (b) Set individual consumer lending targets in all Recovery branches, (c) Evaluate hiring only after 90-day digital migration results, (d) Use freed capacity for inactive PI client activation campaigns.

KEY MANAGEMENT DECISIONS REQUIRED:
- Approve 90-day digital migration program for 12 cash-heavy branches
- Set branch-level and individual consumer lending targets linked to M1 market-share recovery
- Review role allocation in branches flagged for redesign
- Defer all hiring decisions until post-migration capacity is measured
- Link M2 POS deployment to branches where micro staff capacity can be released through cash reduction

LINKAGE:
- M1 Market Share: Recovery branches with low PI penetration get higher acquisition and lending targets
- M2 Product Mix: Branches flagged for consumer lending push (M2) align with sales discipline recommendation
- M2 POS Strategy: Zero-POS branches overlap with branches needing micro/business segment development""")

# ============================================================
# SAVE
# ============================================================
out = 'Module3_Staff_Efficiency.pptx'
if os.path.exists(out):
    out = 'Module3_Staff_Efficiency_v2.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Size: {os.path.getsize(out):,} bytes')
print(f'Slides: {len(prs.slides)}')
