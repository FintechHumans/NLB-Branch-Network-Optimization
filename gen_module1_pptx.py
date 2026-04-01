#!/usr/bin/env python3
"""
NLB Prishtina — Module 1: Market Share
3-Slide Executive PPTX | McKinsey-style | Board-ready
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# DESIGN SYSTEM
# ============================================================
NAVY = RGBColor(0x00, 0x2B, 0x49)
STEEL = RGBColor(0x4A, 0x6F, 0xA5)
CHARCOAL = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============================================================
# DATA
# ============================================================
with open('branch_data_unified.json') as f:
    data = json.load(f)
branches = data['branches']

def s(v, d=0):
    if v is None: return d
    try: return float(v)
    except: return d

# Compute groups
benchmarked = []
for ou, b in branches.items():
    if not b.get('is_benchmarked'):
        continue
    pen = b.get('pop_penetration')
    if pen is None:
        continue
    pen_pct = pen * 100
    biz_pen = b.get('biz_penetration')
    biz_pct = biz_pen * 100 if biz_pen else 0
    benchmarked.append({
        'ou': ou, 'name': b['branch_name'], 'city': b.get('city',''),
        'emp': b['employees'], 'pop': s(b.get('population')),
        'active': s(b['active_clients_pi']), 'inactive': s(b['non_active_clients_pi']),
        'pen': pen_pct, 'biz_pen': biz_pct,
        'act_ratio': s(b['active_clients_pi']) / (s(b['active_clients_pi']) + s(b['non_active_clients_pi'])) * 100 if (s(b['active_clients_pi']) + s(b['non_active_clients_pi'])) > 0 else 0,
    })

benchmarked.sort(key=lambda x: x['pen'], reverse=True)
PI_MEDIAN = 13.1
PI_TARGET = 18.0

defend = [x for x in benchmarked if x['pen'] >= PI_TARGET]
accel = [x for x in benchmarked if PI_MEDIAN <= x['pen'] < PI_TARGET]
recovery = [x for x in benchmarked if x['pen'] < PI_MEDIAN]

# ============================================================
# PRESENTATION SETUP
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank

def add_shape(slide, left, top, width, height, fill=None, border=None, border_width=Pt(1)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = border_width
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, left, top, width, height, text, size=11, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_run_text(tf, text, size=11, bold=False, color=CHARCOAL):
    """Add a new paragraph with a run to an existing text frame."""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

# ============================================================
# SLIDE 1 — MARKET POSITION
# ============================================================
sl = prs.slides.add_slide(blank_layout)

# Navy header bar
add_shape(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.9), fill=NAVY)
add_text(sl, Inches(0.5), Inches(0.12), Inches(8), Inches(0.4),
         'MODULE 1 — MARKET SHARE', size=20, bold=True, color=WHITE)
add_text(sl, Inches(0.5), Inches(0.52), Inches(10), Inches(0.3),
         'NLB Prishtina | Branch Network Transformation | March 2026', size=10, color=RGBColor(0xA0,0xBC,0xD8))

# Executive message bar
shp = add_shape(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.75), fill=RGBColor(0xF0,0xF4,0xF8), border=BORDER)
tf = shp.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "EXECUTIVE INSIGHT:  "
run.font.size = Pt(10)
run.font.bold = True
run.font.color.rgb = NAVY
run2 = p.add_run()
run2.text = "NLB Prishtina captures 15.7% of Kosovo\u2019s addressable population (239,351 active PI clients / 1.52M). Only one main Filiali exceeds the 18% threshold individually: Filiali Peje at 18.2%. All other major Filiali remain below target. 84,765 inactive PI clients (26.2%) represent the fastest reactivation path."
run2.font.size = Pt(10)
run2.font.color.rgb = CHARCOAL

# LEFT COLUMN — Key metrics
left_x = Inches(0.4)
y = Inches(2.1)

# KPI boxes
kpis = [
    ('NETWORK PI PENETRATION', '15.7%', '239,351 active / 1.52M unique pop.', NAVY),
    ('INACTIVE PI CLIENT POOL', '84,765', '26.2% of total PI base dormant', RED),
    ('BRANCHES ABOVE 18% TARGET', '7 of 31', 'Only Filiali Peje among main branches', GREEN),
    ('BUSINESS PENETRATION', '~8.5%', '~12,600 active / ~148,700 businesses', AMBER),
]

for i, (label, value, ctx, accent) in enumerate(kpis):
    ky = y + Inches(i * 1.0)
    box = add_shape(sl, left_x, ky, Inches(3.8), Inches(0.85), fill=WHITE, border=BORDER)
    # Accent left border
    accent_bar = add_shape(sl, left_x, ky, Inches(0.06), Inches(0.85), fill=accent)
    add_text(sl, left_x + Inches(0.2), ky + Inches(0.08), Inches(3.4), Inches(0.2),
             label, size=8, bold=True, color=MID_GRAY)
    add_text(sl, left_x + Inches(0.2), ky + Inches(0.28), Inches(3.4), Inches(0.3),
             value, size=22, bold=True, color=accent)
    add_text(sl, left_x + Inches(0.2), ky + Inches(0.62), Inches(3.4), Inches(0.2),
             ctx, size=8, color=MID_GRAY)

# RIGHT COLUMN — Branch penetration table (top 15 + bottom 5)
rx = Inches(4.6)
add_text(sl, rx, Inches(2.0), Inches(8), Inches(0.3),
         'PI RETAIL PENETRATION BY BRANCH (sorted descending)', size=10, bold=True, color=NAVY)

# Table
from pptx.util import Inches, Pt
chart_branches = [b for b in benchmarked if b['ou'] != '100703']  # Exclude outlier
rows = chart_branches[:31]  # All 31 minus outlier = 30
row_count = len(rows) + 1  # header + data
cols = ['Branch', 'City', 'Pen%', 'Active', 'Inactive', 'Group']

tbl_shape = sl.shapes.add_table(min(row_count, 18), 6,
    rx, Inches(2.35), Inches(8.3), Inches(4.5))
tbl = tbl_shape.table

# Column widths
tbl.columns[0].width = Inches(2.8)
tbl.columns[1].width = Inches(1.3)
tbl.columns[2].width = Inches(0.8)
tbl.columns[3].width = Inches(1.0)
tbl.columns[4].width = Inches(1.0)
tbl.columns[5].width = Inches(1.4)

# Header
headers = ['Branch', 'City', 'Pen%', 'Active PI', 'Inactive', 'Group']
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE

# Data rows (top 17 to fit slide)
display_rows = rows[:17]
for ri, b in enumerate(display_rows):
    pen = b['pen']
    if pen >= PI_TARGET:
        group = 'Defend'
        gcolor = GREEN
    elif pen >= PI_MEDIAN:
        group = 'Acceleration'
        gcolor = AMBER
    else:
        group = 'Recovery'
        gcolor = RED

    vals = [
        b['name'],
        b['city'],
        f"{pen:.1f}%",
        f"{int(b['active']):,}",
        f"{int(b['inactive']):,}",
        group,
    ]
    for ci, v in enumerate(vals):
        cell = tbl.cell(ri + 1, ci)
        cell.text = v
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.color.rgb = gcolor if ci == 2 or ci == 5 else CHARCOAL
            if ci == 2:
                p.font.bold = True

# Note below table
add_text(sl, rx, Inches(6.95), Inches(8), Inches(0.4),
         'Mitrovice Veriore excluded (structural outlier, 116% pen). Full 30-branch chart in HTML dashboard. Median: 13.1% | Target: 18.0%',
         size=7, color=MID_GRAY)

# Speaker notes
notes = sl.notes_slide
notes.notes_text_frame.text = """SLIDE 1 — MARKET POSITION

Key points:
- Network PI penetration: 15.7% (239,351 active / 1.52M unique population)
- Only one main Filiali exceeds 18%: Filiali Peje at 18.2%. All other major Filiali below target.
- The 7 defend branches are predominantly small ekspoziture: Hani i Elezit (39.7%), Shterpce (27.5%), Decan (27.1%), Kacanik (21.4%), Gracanice (20.0%), Obiliq (18.1%)
- Peje is the proof point that scale and market dominance can coexist

PRISHTINA CITY-LEVEL CONTEXT:
At city level, Prishtina stands at 21.4% PI penetration (48,619 active clients across 227,466 population). This means the city is NOT underpenetrated overall; the issue is internal branch allocation and outlet role clarity, not market weakness. 4 active outlets + 1 ghost branch (100101, 7 FTE, zero output) share this market.

Mitrovice Veriore excluded from chart — 116% penetration due to cross-municipality registration. Treated as structural outlier.
"""

# ============================================================
# SLIDE 2 — ACTION GROUPS & REACTIVATION
# ============================================================
sl2 = prs.slides.add_slide(blank_layout)

# Header
add_shape(sl2, Inches(0), Inches(0), SLIDE_W, Inches(0.9), fill=NAVY)
add_text(sl2, Inches(0.5), Inches(0.12), Inches(10), Inches(0.4),
         'BRANCH ACTION GROUPS \u2014 PI RETAIL PENETRATION', size=20, bold=True, color=WHITE)
add_text(sl2, Inches(0.5), Inches(0.52), Inches(10), Inches(0.3),
         'Recovery / Acceleration / Defend classification based on median (13.1%) and target (18.0%)', size=10, color=RGBColor(0xA0,0xBC,0xD8))

# Three group cards
groups = [
    ('DEFEND / DEEPEN', f'{len(defend)} branches', '\u226518% target', GREEN, RGBColor(0xEC,0xFD,0xF5),
     'Protect base. Deepen cross-sell.\nOnly Filiali Peje (18.2%) is a main branch.\nOthers: Hani Elezit, Shterpce, Decan,\nKacanik, Gracanice, Obiliq.'),
    ('ACCELERATION', f'{len(accel)} branches', '13.1%\u201318%', AMBER, RGBColor(0xFF,0xFB,0xEB),
     f'Highest-ROI group. Gap: +15,262 clients.\nPriorities: Gjilan (+3,392), Mitrovice\n(+3,156), Besiane (+2,278).\nKline: 62.6% active ratio \u2014 pilot\nfor reactivation program.'),
    ('RECOVERY', f'{len(recovery)} branches', '<13.1% median', RED, RGBColor(0xFE,0xF2,0xF2),
     f'Nearly half the network. 40,543 inactive.\nCore Filiali: Prishtine (7.3%), Prizren\n(7.9%), Ferizaj (12.4%), Gjakove (11.2%).\nPrizren: worst major city at 10.9% combined.'),
]

for i, (title, count, threshold, color, bg, desc) in enumerate(groups):
    gx = Inches(0.4 + i * 4.2)
    gy = Inches(1.15)
    box = add_shape(sl2, gx, gy, Inches(3.9), Inches(2.6), fill=bg, border=color, border_width=Pt(1.5))
    add_text(sl2, gx + Inches(0.2), gy + Inches(0.12), Inches(3.5), Inches(0.25),
             title, size=11, bold=True, color=color)
    add_text(sl2, gx + Inches(0.2), gy + Inches(0.4), Inches(3.5), Inches(0.4),
             count, size=28, bold=True, color=color)
    add_text(sl2, gx + Inches(0.2), gy + Inches(0.85), Inches(3.5), Inches(0.2),
             threshold, size=9, color=MID_GRAY)
    add_text(sl2, gx + Inches(0.2), gy + Inches(1.15), Inches(3.5), Inches(1.3),
             desc, size=8, color=CHARCOAL)

# Reactivation scenario box
shp = add_shape(sl2, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.85), fill=RGBColor(0xF0,0xF4,0xF8), border=STEEL)
tf = shp.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "REACTIVATION SCENARIO:  "
run.font.size = Pt(10)
run.font.bold = True
run.font.color.rgb = NAVY
run2 = p.add_run()
run2.text = "The network holds 84,765 inactive PI clients. Reactivating 25% adds 21,191 clients \u2192 pen rises from 15.7% to 17.1%. Reactivating 50% \u2192 18.5%. Top pools: 1001 Prishtine (6,976), 0001 Prishtine (5,899), Peje (5,465), Prizren (4,514), Ferizaj (4,386)."
run2.font.size = Pt(9)
run2.font.color.rgb = CHARCOAL

# Priority action table
add_text(sl2, Inches(0.4), Inches(5.05), Inches(8), Inches(0.3),
         'PRIORITY ACTIONS \u2014 IMMEDIATE', size=11, bold=True, color=NAVY)

actions = [
    ['Suhareke (100502)', '12.8%', '+165 to median', 'Quick win \u2014 one good month closes gap'],
    ['Ferizaj (1003)', '12.4%', '+782 to median', '19 FTE \u2014 focused campaign closes in one quarter'],
    ['Gjakove (1004)', '11.2%', '+1,535 to median', '3,850 inactive \u2014 reactivation exceeds gap'],
    ['Kline (100602)', '14.1%', '+1,181 to 18%', '62.6% active ratio \u2014 reactivation pilot'],
    ['Gjilan (1002)', '13.9%', '+3,392 to 18%', 'Main Filiali, 12 FTE \u2014 dedicated campaign'],
    ['Prizren city', '10.9%', 'City strategy', 'Worst major city \u2014 needs city-level plan'],
]

atbl = sl2.shapes.add_table(len(actions) + 1, 4,
    Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.8))
at = atbl.table
at.columns[0].width = Inches(2.5)
at.columns[1].width = Inches(1.2)
at.columns[2].width = Inches(2.0)
at.columns[3].width = Inches(6.8)

for ci, h in enumerate(['Branch / City', 'Current Pen%', 'Gap', 'Management Action']):
    cell = at.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE

for ri, row in enumerate(actions):
    for ci, v in enumerate(row):
        cell = at.cell(ri + 1, ci)
        cell.text = v
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.color.rgb = CHARCOAL

# Prishtina callout
shp = add_shape(sl2, Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.0), fill=None)  # spacer

# Speaker notes
notes2 = sl2.notes_slide
notes2.notes_text_frame.text = """SLIDE 2 — ACTION GROUPS

PRISHTINA CITY-LEVEL CONTEXT:
At city level, Prishtina stands at 21.4% PI penetration. The city is NOT underpenetrated. The issue is internal branch allocation and outlet role clarity, not market weakness. The ghost branch (100101, 7 FTE, zero output) absorbs 2.6% of total network FTE.

KLINE REACTIVATION PILOT:
Kline has the worst active ratio in the network (62.6%). 2,574 inactive vs 4,309 active. If 50% reactivated, pen jumps from 14.1% to ~18.3%. Recommend as the pilot branch for network-wide reactivation program.

PRIZREN CITY STRATEGY:
At 10.9% combined (2 branches, 16 FTE, 147K population), Prizren is the most underpenetrated major city. This is not a branch problem — it requires a city-level sales strategy.

METHODOLOGY: In multi-branch cities (Prishtina, Prizren, Mitrovice), branch-level pen% is computed against full city population. City-level pen is the correct cross-market comparison. Branch-level is for internal allocation only.
"""

# ============================================================
# SLIDE 3 — BUSINESS OPPORTUNITY
# ============================================================
sl3 = prs.slides.add_slide(blank_layout)

# Header
add_shape(sl3, Inches(0), Inches(0), SLIDE_W, Inches(0.9), fill=NAVY)
add_text(sl3, Inches(0.5), Inches(0.12), Inches(10), Inches(0.4),
         'BUSINESS MARKET \u2014 THE LARGEST WHITE SPACE', size=20, bold=True, color=WHITE)
add_text(sl3, Inches(0.5), Inches(0.52), Inches(10), Inches(0.3),
         'Business penetration at ~8.5% vs retail 15.7% \u2014 proportionally the bigger structural gap', size=10, color=RGBColor(0xA0,0xBC,0xD8))

# Executive bar
shp = add_shape(sl3, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.65), fill=RGBColor(0xF0,0xF4,0xF8), border=BORDER)
tf = shp.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "KEY FINDING:  "
run.font.size = Pt(10)
run.font.bold = True
run.font.color.rgb = NAVY
run2 = p.add_run()
run2.text = "NLB Prishtina captures only ~8.5% of ~148,700 registered businesses. Every major Filiali except those in smaller cities is below the 18% target. Business active ratio (~58.6%) is materially worse than PI (73.8%), signaling higher churn/dormancy in the business segment."
run2.font.size = Pt(9)
run2.font.color.rgb = CHARCOAL

# KPI cards
biz_kpis = [
    ('BUSINESS PENETRATION', '~8.5%', '~12,600 active / ~148,700 businesses', AMBER),
    ('BUSINESS ACTIVE RATIO', '~58.6%', 'vs PI active ratio of 73.8%', RED),
    ('BIZ BRANCHES \u226518%', '11 of 31', 'Mostly small ekspoziture', GREEN),
]

for i, (label, value, ctx, accent) in enumerate(biz_kpis):
    bx = Inches(0.4 + i * 4.2)
    by = Inches(2.0)
    box = add_shape(sl3, bx, by, Inches(3.9), Inches(0.85), fill=WHITE, border=BORDER)
    accent_bar = add_shape(sl3, bx, by, Inches(0.06), Inches(0.85), fill=accent)
    add_text(sl3, bx + Inches(0.2), by + Inches(0.08), Inches(3.5), Inches(0.2),
             label, size=8, bold=True, color=MID_GRAY)
    add_text(sl3, bx + Inches(0.2), by + Inches(0.3), Inches(3.5), Inches(0.3),
             value, size=22, bold=True, color=accent)
    add_text(sl3, bx + Inches(0.2), by + Inches(0.62), Inches(3.5), Inches(0.2),
             ctx, size=8, color=MID_GRAY)

# Business penetration table — top under-penetrated markets
add_text(sl3, Inches(0.4), Inches(3.1), Inches(8), Inches(0.3),
         'BUSINESS PENETRATION \u2014 KEY MARKETS (sorted by opportunity size)', size=11, bold=True, color=NAVY)

# Sort by businesses descending (largest markets)
biz_sorted = sorted(benchmarked, key=lambda x: x.get('pop', 0), reverse=True)
biz_rows = [(b['name'], b['city'], f"{b['biz_pen']:.1f}%",
             f"{int(b.get('pop',0)):,}", 'Defend' if b['biz_pen'] >= 18 else ('Accel' if b['biz_pen'] >= 15.2 else 'Recovery'))
            for b in biz_sorted[:15]]

btbl = sl3.shapes.add_table(len(biz_rows) + 1, 5,
    Inches(0.4), Inches(3.45), Inches(12.5), Inches(3.5))
bt = btbl.table
bt.columns[0].width = Inches(3.2)
bt.columns[1].width = Inches(1.8)
bt.columns[2].width = Inches(1.5)
bt.columns[3].width = Inches(2.0)
bt.columns[4].width = Inches(4.0)

for ci, h in enumerate(['Branch', 'City', 'Biz Pen%', 'Population', 'Group']):
    cell = bt.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE

for ri, row in enumerate(biz_rows):
    for ci, v in enumerate(row):
        cell = bt.cell(ri + 1, ci)
        cell.text = v
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            if ci == 2:
                p.font.bold = True
                pen_val = float(v.replace('%',''))
                p.font.color.rgb = GREEN if pen_val >= 18 else (AMBER if pen_val >= 15.2 else RED)
            elif ci == 4:
                p.font.color.rgb = GREEN if 'Defend' in v else (AMBER if 'Accel' in v else RED)
            else:
                p.font.color.rgb = CHARCOAL

# Data quality note
add_text(sl3, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
         'Note: Active business account counts are estimated (penetration ratio \u00d7 registered businesses). Inactive business counts derived from business active ratio. Treat as directional.',
         size=7, color=MID_GRAY)

# Speaker notes
notes3 = sl3.notes_slide
notes3.notes_text_frame.text = """SLIDE 3 — BUSINESS OPPORTUNITY

Business penetration is structurally weaker than retail:
- Network biz pen: ~8.5% vs PI pen 15.7%
- Business active ratio: ~58.6% vs PI 73.8%
- Every major Filiali is below 18% in business penetration

Largest opportunities by market size:
- Prishtina: 21,396 businesses, 0001 at 10.5%, 1001 at 6.7%
- Prizren: 7,119 businesses at 9.1%
- Ferizaj: 6,371 businesses at 11.4%
- Gjilan: 4,317 businesses at 10.6%
- Gjakove: 4,005 businesses at 10.9%
- Peje: 4,172 businesses at 16.4%

DATA QUALITY:
- Active business counts are ESTIMATED from penetration ratio, not directly counted
- Inactive business counts are ESTIMATED from business active ratio
- These are directional indicators, not precise counts
- PI active/inactive counts are direct from source — high confidence
"""

# ============================================================
# SAVE
# ============================================================
fname = 'Module1_Market_Share.pptx'
try:
    prs.save(fname)
except PermissionError:
    fname = 'Module1_Market_Share_v2.pptx'
    prs.save(fname)

import os
print(f"Saved: {fname}")
print(f"Size: {os.path.getsize(fname):,} bytes")
print(f"Slides: {len(prs.slides)}")
